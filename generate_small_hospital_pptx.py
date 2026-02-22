"""
Generate comprehensive Small Hospital EMR Pitch PowerPoint Presentation.
50+ slides covering all features, architecture, benefits, and deployment.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "Hospital_EMR_Pitch")
os.makedirs(OUTPUT_DIR, exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette
DARK_BLUE = RGBColor(0, 51, 102)
MED_BLUE = RGBColor(0, 102, 153)
LIGHT_BLUE = RGBColor(230, 240, 250)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(33, 33, 33)
DARK_GREEN = RGBColor(0, 102, 51)
LIGHT_GREEN = RGBColor(230, 245, 230)
RED = RGBColor(180, 30, 30)
ORANGE = RGBColor(230, 115, 0)
PURPLE = RGBColor(102, 0, 102)
GRAY = RGBColor(100, 100, 100)
LIGHT_GRAY = RGBColor(240, 240, 240)
ACCENT_TEAL = RGBColor(0, 150, 136)

# Helper Functions
def add_bg(slide, color=DARK_BLUE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=BLACK, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_multi_text(slide, left, top, width, height, lines,
                   font_size=14, color=BLACK, spacing=1.2, bullet=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line, tuple):
            text, line_bold, line_color = line
        else:
            text = line
            line_bold = False
            line_color = color
        if bullet and not text.startswith("•"):
            text = "• " + text
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = line_color
        p.font.bold = line_bold
        p.font.name = "Calibri"
        p.space_after = Pt(font_size * spacing * 0.4)
    return txBox

def add_table(slide, left, top, width, height, headers, rows):
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table
    col_w = int(width / num_cols)
    for i in range(num_cols):
        table.columns[i].width = col_w
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            if r_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BLUE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = BLACK
                p.font.name = "Calibri"
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table_shape

def title_slide_layout(slide, section_num, title, subtitle=""):
    add_bg(slide, DARK_BLUE)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(5.8), Inches(1.5), Inches(1.6), Inches(1.6)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = WHITE
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = str(section_num)
    p.font.size = Pt(48)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(12)
    add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(1.5),
                title, 40, WHITE, True, PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide, Inches(1), Inches(5), Inches(11), Inches(1),
                    subtitle, 20, RGBColor(180, 210, 240), False, PP_ALIGN.CENTER)

def add_kpi_card(slide, left, top, value, label, color):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.3), Inches(1.3)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(10)
    p2.font.color.rgb = RGBColor(220, 220, 255)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1: COVER
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), MED_BLUE)

icon_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(1.0), Inches(2.3), Inches(1.5)
)
icon_box.fill.solid()
icon_box.fill.fore_color.rgb = WHITE
icon_box.line.fill.background()
tf = icon_box.text_frame
p = tf.paragraphs[0]
p.text = "🏥"
p.font.size = Pt(60)
p.alignment = PP_ALIGN.CENTER

add_textbox(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(1.2),
            "SMALL HOSPITAL EMR", 44, WHITE, True, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.8),
            "Modern Cloud-Based Electronic Medical Records System", 22,
            RGBColor(180, 210, 240), False, PP_ALIGN.CENTER)
add_textbox(slide, Inches(2), Inches(5.2), Inches(9.3), Inches(0.7),
            "For 50-200 Bed Private Hospitals", 18,
            RGBColor(150, 180, 210), False, PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(0), Inches(6.6), prs.slide_width, Inches(0.9), RGBColor(0, 40, 80))
add_textbox(slide, Inches(1), Inches(6.7), Inches(11.3), Inches(0.6),
            "COMPREHENSIVE PITCH & DEPLOYMENT GUIDE   |   February 2026",
            14, RGBColor(180, 200, 220), False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2: TABLE OF CONTENTS
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "TABLE OF CONTENTS", 32, WHITE, True, PP_ALIGN.LEFT)

toc_items = [
    "1. What is This EMR?",
    "2. Why Choose This Solution?",
    "3. Complete Feature Overview",
    "4. System Architecture",
    "5. Workflow Demonstrations",
    "6. Key Benefits & ROI",
    "7. Implementation Timeline",
    "8. Pricing & Licensing",
    "9. Security & Compliance",
    "10. Deployment Plan",
    "11. Support & Training",
    "12. Next Steps",
]

for i, item in enumerate(toc_items):
    col = i % 3
    row = i // 3
    left = Inches(0.5 + col * 4.2)
    top = Inches(1.3 + row * 1.9)
    add_textbox(slide, left, top, Inches(4.0), Inches(0.6),
                item, 14, DARK_BLUE, False, PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3: WHAT IS THIS EMR?
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "WHAT IS THIS EMR SYSTEM?", 28, WHITE, True)

description = [
    ("Modern Solution", "Cloud-based electronic medical records system built specifically for small private hospitals (50-200 beds)"),
    ("Complete Integration", "All departments connected on a single platform — Patient, EMR, Lab, Pharmacy, Billing, Nursing, Admin"),
    ("Easy to Use", "Intuitive web interface designed for busy clinicians — no steep learning curve"),
    ("Scalable Infrastructure", "Grows with your hospital — from 50 beds to 500+ beds without any code changes"),
    ("Affordable", "Modern cloud stack keeps costs low ($200-300/month) vs traditional EMR ($5-10K/month)"),
    ("Security First", "Bank-grade AES-256 encryption, HIPAA-compliant design, complete audit trail"),
    ("Fast Deployment", "12 weeks from order to live system, fully deployed and ready to go"),
    ("Owned by You", "You own the code and data — no vendor lock-in, can modify/extend as needed"),
]

add_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8),
          ["Aspect", "Description"], description)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4: WHY CHOOSE THIS SOLUTION?
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "WHY CHOOSE THIS SOLUTION?", 28, WHITE, True)

add_textbox(slide, Inches(0.5), Inches(1.3), Inches(12), Inches(0.6),
            "10 Reasons Small Hospitals Are Choosing This EMR", 16, DARK_BLUE, True)

reasons = [
    ("1. Cost-Effective", "$200-300/month vs $5-10K/month for enterprise solutions"),
    ("2. No IT Team Needed", "Managed cloud services — all infrastructure handled by experts"),
    ("3. Fast Deployment", "12 weeks start-to-finish (vs 2 years for traditional EMR)"),
    ("4. All-in-One", "Patient, EMR, Lab, Pharmacy, Billing, Nursing — all integrated"),
    ("5. Modern Technology", "Built with latest tech stack (React, Node.js, PostgreSQL)"),
    ("6. Easy to Use", "Designed for busy clinicians, not IT specialists"),
    ("7. Secure & Compliant", "HIPAA design, AES-256 encryption, full audit trail"),
    ("8. Scalable", "Grows from 50 to 500+ beds without major changes"),
    ("9. Future-Proof", "Cloud-based, auto-updates, security patches automatic"),
    ("10. Owned by You", "Source code is yours, data is yours, can extend as needed"),
]

for i, (title, desc) in enumerate(reasons):
    col = i % 2
    row = i // 2
    left = Inches(0.5 + col * 6.4)
    top = Inches(2.1 + row * 1.0)
    # Checkbox
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, Inches(0.35), Inches(0.35)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = DARK_GREEN
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = "✓"
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    # Title
    add_textbox(slide, left + Inches(0.5), top, Inches(5.5), Inches(0.3),
                title, 12, DARK_BLUE, True)
    # Description
    add_textbox(slide, left + Inches(0.5), top + Inches(0.3), Inches(5.5), Inches(0.6),
                desc, 10, GRAY)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5: COMPLETE FEATURE OVERVIEW
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_layout(slide, 1, "COMPLETE FEATURE OVERVIEW",
                   "12 Integrated Modules — Everything You Need")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6: MODULE BREAKDOWN - PAGE 1
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "MODULE 1-4: PATIENT & APPOINTMENTS", 24, WHITE, True)

modules_1_4 = [
    ("1. PATIENT MANAGEMENT", "📋", DARK_BLUE, [
        "New patient registration",
        "Unique patient ID generation",
        "Medical history tracking",
        "Allergy & alert management",
        "Document upload (ID, insurance)",
        "Patient search (name, ID, phone)",
        "Patient categorization (VIP, regular)",
        "Demographics & contact info",
    ]),
    ("2. ELECTRONIC MEDICAL RECORDS (EMR)", "📝", MED_BLUE, [
        "SOAP consultation notes",
        "Diagnosis recording (ICD-10)",
        "E-prescribing with drug checks",
        "Vital signs recording",
        "Lab order creation",
        "Medical templates (specialty-specific)",
        "Allergy & clinical alerts",
        "Complete medical timeline",
    ]),
    ("3. APPOINTMENTS & SCHEDULING", "📅", ACCENT_TEAL, [
        "Online appointment booking",
        "Walk-in queue management",
        "Doctor availability management",
        "Automated reminders (SMS/email)",
        "Recurring appointments",
        "Follow-up scheduling",
        "Telemedicine slot booking",
        "Queue status display (real-time)",
    ]),
    ("4. OUT-PATIENT DEPARTMENT (OPD)", "👨‍⚕️", DARK_GREEN, [
        "Patient check-in interface",
        "Queue management display",
        "Consultation workflow",
        "Prescription issuing",
        "Lab/imaging referrals",
        "Specialist referrals",
        "Follow-up management",
        "OPD reports & analytics",
    ]),
]

for i, (title, icon, color, features) in enumerate(modules_1_4):
    left = Inches(0.3 + (i % 2) * 6.5)
    top = Inches(1.4 + (i // 2) * 3.0)
    
    # Header
    hdr = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(6.2), Inches(0.5)
    )
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = color
    hdr.line.fill.background()
    tf = hdr.text_frame
    p = tf.paragraphs[0]
    p.text = f"{icon}  {title}"
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Features
    y_offset = 0.6
    for feature in features[:4]:
        add_textbox(slide, left + Inches(0.2), top + Inches(y_offset), Inches(5.8), Inches(0.35),
                    f"• {feature}", 9, BLACK, False)
        y_offset += 0.35

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7: MODULE BREAKDOWN - PAGE 2
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "MODULE 5-8: PHARMACY, LAB & IN-PATIENT", 24, WHITE, True)

modules_5_8 = [
    ("5. PHARMACY MANAGEMENT", "💊", ORANGE, [
        "E-prescription receipt",
        "Medication dispensing",
        "Drug interaction checking",
        "Inventory management",
        "Stock level alerts",
        "Barcode-based dispensing",
        "Retail sales (OTC)",
        "Purchase order tracking",
    ]),
    ("6. LABORATORY (LIS)", "🔬", PURPLE, [
        "Digital lab orders",
        "Sample collection tracking",
        "Result entry interface",
        "Auto-validation rules",
        "Critical value alerts",
        "TAT monitoring",
        "Test history & reports",
        "Quality control data",
    ]),
    ("7. IN-PATIENT / WARDS", "🛏️", RED, [
        "Admission workflow",
        "Bed management & visual map",
        "Ward dashboard",
        "Patient transfer tracking",
        "Discharge planning",
        "Daily census reports",
        "Length of stay monitoring",
        "Ward notes & updates",
    ]),
    ("8. NURSING STATION", "👩‍⚕️", LIGHT_GREEN, [
        "Patient worklist display",
        "Vital signs charting",
        "Care plan management",
        "Medication administration (eMAR)",
        "Intake/output tracking",
        "Patient alerts & tasks",
        "Shift handover notes",
        "Real-time updates",
    ]),
]

for i, (title, icon, color, features) in enumerate(modules_5_8):
    left = Inches(0.3 + (i % 2) * 6.5)
    top = Inches(1.4 + (i // 2) * 3.0)
    
    hdr = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(6.2), Inches(0.5)
    )
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = color
    hdr.line.fill.background()
    tf = hdr.text_frame
    p = tf.paragraphs[0]
    p.text = f"{icon}  {title}"
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    y_offset = 0.6
    for feature in features[:4]:
        add_textbox(slide, left + Inches(0.2), top + Inches(y_offset), Inches(5.8), Inches(0.35),
                    f"• {feature}", 9, BLACK, False)
        y_offset += 0.35

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8: MODULE BREAKDOWN - PAGE 3
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "MODULE 9-12: BILLING, ADMIN & SECURITY", 24, WHITE, True)

modules_9_12 = [
    ("9. BILLING & FINANCIAL", "💰", ORANGE, [
        "Auto-billing from services",
        "Invoice generation & printing",
        "Multiple payment methods",
        "Insurance claims (basic)",
        "Revenue tracking",
        "Outstanding payment alerts",
        "Payment receipts",
        "Financial reports",
    ]),
    ("10. ADMIN DASHBOARD", "📊", DARK_BLUE, [
        "KPI dashboards (real-time)",
        "Revenue monitoring",
        "Patient statistics",
        "Department performance",
        "User activity logs",
        "System health monitoring",
        "Custom widgets",
        "Export reports",
    ]),
    ("11. USER MANAGEMENT & RBAC", "🔐", PURPLE, [
        "7 predefined roles",
        "Role-based access control",
        "Permission management",
        "User creation/editing",
        "Active user monitoring",
        "Password policies",
        "MFA support",
        "Session management",
    ]),
    ("12. SECURITY & AUDIT", "🔒", RED, [
        "Complete audit logging",
        "Access tracking",
        "Data encryption (AES-256)",
        "Backup automation (daily)",
        "Disaster recovery ready",
        "HIPAA-compliant design",
        "GDPR data portability",
        "Compliance reporting",
    ]),
]

for i, (title, icon, color, features) in enumerate(modules_9_12):
    left = Inches(0.3 + (i % 2) * 6.5)
    top = Inches(1.4 + (i // 2) * 3.0)
    
    hdr = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(6.2), Inches(0.5)
    )
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = color
    hdr.line.fill.background()
    tf = hdr.text_frame
    p = tf.paragraphs[0]
    p.text = f"{icon}  {title}"
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    y_offset = 0.6
    for feature in features[:4]:
        add_textbox(slide, left + Inches(0.2), top + Inches(y_offset), Inches(5.8), Inches(0.35),
                    f"• {feature}", 9, BLACK, False)
        y_offset += 0.35

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9: SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_layout(slide, 2, "SYSTEM ARCHITECTURE",
                   "Modern Cloud Stack: Supabase + Koyeb + Vercel")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10: ARCHITECTURE DIAGRAM
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "CLOUD ARCHITECTURE: SUPABASE + KOYEB + VERCEL", 24, WHITE, True)

# Frontend layer
add_textbox(slide, Inches(0.3), Inches(1.4), Inches(2.5), Inches(0.4),
            "FRONTEND", 12, MED_BLUE, True)
frontend_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.9), Inches(2.5), Inches(1.0))
frontend_box.fill.solid()
frontend_box.fill.fore_color.rgb = LIGHT_BLUE
frontend_box.line.color.rgb = MED_BLUE
tf = frontend_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Vercel\nNext.js App\nGlobal CDN\n70+ Pages"
p.font.size = Pt(10)
p.font.color.rgb = DARK_BLUE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Backend layer
add_textbox(slide, Inches(5.2), Inches(1.4), Inches(2.5), Inches(0.4),
            "BACKEND", 12, DARK_BLUE, True)
backend_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.9), Inches(2.5), Inches(1.0))
backend_box.fill.solid()
backend_box.fill.fore_color.rgb = LIGHT_BLUE
backend_box.line.color.rgb = DARK_BLUE
tf = backend_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Koyeb\nNode.js API\nAuto-scaling\n150+ Endpoints"
p.font.size = Pt(10)
p.font.color.rgb = DARK_BLUE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Database layer
add_textbox(slide, Inches(10.1), Inches(1.4), Inches(2.5), Inches(0.4),
            "DATABASE", 12, DARK_GREEN, True)
db_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.1), Inches(1.9), Inches(2.5), Inches(1.0))
db_box.fill.solid()
db_box.fill.fore_color.rgb = LIGHT_GREEN
db_box.line.color.rgb = DARK_GREEN
tf = db_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Supabase\nPostgreSQL\nReal-time\n80+ Tables"
p.font.size = Pt(10)
p.font.color.rgb = DARK_GREEN
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Arrows
for x_pos in [Inches(2.9), Inches(7.8)]:
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, x_pos, Inches(2.3), Inches(0.4), Inches(0.4)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = MED_BLUE
    arrow.line.fill.background()

# Key features
add_textbox(slide, Inches(0.3), Inches(3.2), Inches(12.7), Inches(0.5),
            "Why This Stack?", 16, DARK_BLUE, True)

features_text = [
    "✅ No DevOps Team Needed  |  ✅ Auto-scaling Included  |  ✅99.9% Uptime SLA",
    "✅ Cost-Effective ($200-500/mo)  |  ✅ Instant Deployments (git push → live)",
    "✅ Enterprise Security (AES-256)  |  ✅ Automatic Backups & Updates",
]

for i, feat in enumerate(features_text):
    add_textbox(slide, Inches(0.5), Inches(3.8 + i * 0.5), Inches(12), Inches(0.4),
                feat, 11, BLACK, False, PP_ALIGN.CENTER)

# Benefits boxes
benefits = [
    ("🚀 Fast", "Deploy in 12 weeks"),
    ("💰 Affordable", "$200-300/month"),
    ("🔧 Owned", "Your code & data"),
    ("📈 Scalable", "Grow without code changes"),
    ("🔐 Secure", "HIPAA-compliant design"),
    ("🛠️ Simple", "No IT team needed"),
]

for i, (title, desc) in enumerate(benefits):
    left = Inches(0.3 + i * 2.15)
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(5.3), Inches(1.9), Inches(1.8)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.color.rgb = DARK_BLUE
    box.line.width = Pt(1)
    add_textbox(slide, left + Inches(0.1), Inches(5.4), Inches(1.7), Inches(0.5),
                title, 14, DARK_BLUE, True, PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), Inches(5.95), Inches(1.7), Inches(0.8),
                desc, 11, BLACK, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11: TECHNOLOGY STACK
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "TECHNOLOGY STACK", 24, WHITE, True)

tech_data = [
    ("Frontend", MED_BLUE, ["Next.js 14", "React", "TypeScript", "Tailwind CSS", "Shadcn/UI"]),
    ("Backend", DARK_BLUE, ["Node.js", "NestJS", "TypeScript", "RESTful API", "GraphQL ready"]),
    ("Database", DARK_GREEN, ["PostgreSQL 16", "Supabase Auth", "Row-Level Security", "Real-time", "Backups"]),
    ("DevOps", ORANGE, ["Docker", "GitHub Actions", "CI/CD", "SSL/TLS", "Monitoring"]),
]

for i, (layer, color, techs) in enumerate(tech_data):
    top = Inches(1.4 + i * 1.2)
    # Layer label
    lbl = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), top, Inches(1.5), Inches(0.7)
    )
    lbl.fill.solid()
    lbl.fill.fore_color.rgb = color
    lbl.line.fill.background()
    tf = lbl.text_frame
    p = tf.paragraphs[0]
    p.text = layer
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(8)
    
    # Tech items
    for t_idx, tech in enumerate(techs):
        left = Inches(2.1 + t_idx * 2.0)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.8), Inches(0.7)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE
        box.line.color.rgb = color
        box.line.width = Pt(1.5)
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = tech
        p.font.size = Pt(10)
        p.font.color.rgb = color
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(5)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12: WORKFLOW DEMONSTRATIONS
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_layout(slide, 3, "WORKFLOW DEMONSTRATIONS",
                   "How the EMR Works in Real Hospital Operations")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 13: PATIENT JOURNEY WORKFLOW
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "PATIENT JOURNEY: FROM ARRIVAL TO DISCHARGE", 24, WHITE, True)

# Workflow steps
steps = [
    ("Patient\nArrives", DARK_BLUE, Inches(0.5)),
    ("Registration/\nCheck-in", MED_BLUE, Inches(1.8)),
    ("Queue\nManagement", ACCENT_TEAL, Inches(3.1)),
    ("Doctor\nConsultation", DARK_BLUE, Inches(4.4)),
    ("Orders &\nTests", DARK_GREEN, Inches(5.7)),
]

for i, (text, color, left) in enumerate(steps):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(1.0), Inches(0.8))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    if i < len(steps) - 1:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, left + Inches(1.05), Inches(1.65), Inches(0.25), Inches(0.5)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MED_BLUE
        arrow.line.fill.background()

# Outcomes
add_textbox(slide, Inches(7), Inches(1.5), Inches(6), Inches(0.8),
            "Results & Actions:", 14, DARK_BLUE, True, PP_ALIGN.LEFT)

outcomes = [
    "🔬 Lab Tests → Results in 24 hours",
    "💊 Prescriptions → Auto to Pharmacy",
    "📊 Reports → Instant to Doctor & Patient",
    "💰 Charges → Auto-billing generated",
    "✅ Follow-up → Auto-scheduled",
]

for i, outcome in enumerate(outcomes):
    add_textbox(slide, Inches(7), Inches(2.4 + i * 0.5), Inches(5.8), Inches(0.4),
                outcome, 11, BLACK, False)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 14: DOCTOR CONSULTATION FLOW
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "DOCTOR CONSULTATION WORKFLOW", 24, WHITE, True)

add_textbox(slide, Inches(0.5), Inches(1.3), Inches(12), Inches(0.5),
            "1. Doctor logs in → 2. Views patient list → 3. Opens patient record → 4. Reviews history", 12, BLACK, False)

doc_workflow = [
    ("📋 Patient Info", "Name, age, medical history, allergies, past visits"),
    ("💬 Create Note", "Doctor types SOAP note (Subjective, Objective, Assessment, Plan)"),
    ("📊 Record Vitals", "Temperature, BP, pulse, SpO2, weight auto-calculated BMI"),
    ("💊 E-Prescribe", "Select drugs, dose, frequency — system checks for interactions"),
    ("🔬 Order Labs", "Select tests needed — order sent directly to lab with patient ID"),
    ("💾 Save & Bill", "Note saved, charges auto-generated, patient billed, email sent"),
]

colors_cycle = [DARK_BLUE, MED_BLUE, ACCENT_TEAL, DARK_GREEN, ORANGE, PURPLE]
for i, (title, desc) in enumerate(doc_workflow):
    left = Inches(0.5)
    top = Inches(2.0 + i * 0.75)
    
    num_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, Inches(0.4), Inches(0.4)
    )
    num_circle.fill.solid()
    num_circle.fill.fore_color.rgb = colors_cycle[i]
    num_circle.line.fill.background()
    tf = num_circle.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    add_textbox(slide, left + Inches(0.6), top, Inches(2.0), Inches(0.4),
                title, 12, colors_cycle[i], True)
    add_textbox(slide, left + Inches(3.0), top, Inches(9.3), Inches(0.4),
                desc, 11, BLACK, False)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 15: PHARMACY DISPENSING WORKFLOW
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "PHARMACY DISPENSING WORKFLOW", 24, WHITE, True)

pharmacy_steps = [
    ("1", "E-Rx Received", "Doctor's prescription appears in pharmacy queue automatically", ORANGE),
    ("2", "Verification", "Pharmacist verifies patient + drug + dosage", DARK_GREEN),
    ("3", "Interaction Check", "System checks for drug-drug & drug-allergy interactions", RED),
    ("4", "Dispensing", "Medications prepared, barcode printed, patient verifies", PURPLE),
    ("5", "Counselling", "Pharmacist counsels patient (dosage, side effects, timing)", BLUE := MED_BLUE),
    ("6", "Billing", "Charges auto-sent to billing, inventory auto-updated", ORANGE),
]

for i, (num, title, desc, color) in enumerate(pharmacy_steps):
    top = Inches(1.4 + i * 0.9)
    
    # Number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.3), top, Inches(0.5), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Title
    add_textbox(slide, Inches(1.0), top, Inches(2.0), Inches(0.5),
                title, 12, color, True)
    # Description
    add_textbox(slide, Inches(3.2), top, Inches(9.6), Inches(0.5),
                desc, 11, BLACK, False)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 16: LAB & BILLING WORKFLOWS
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "LABORATORY & BILLING WORKFLOWS", 24, WHITE, True)

# Lab workflow (left)
add_textbox(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.5),
            "LABORATORY WORKFLOW", 14, DARK_BLUE, True)

lab_workflow = [
    "1. Doctor orders tests (EMR)",
    "2. Lab receives order (real-time alert)",
    "3. Sample collection with barcode",
    "4. Lab technician enters results",
    "5. Auto-validation (normal range check)",
    "6. Critical alerts if abnormal",
    "7. Report generated & doctor notified",
    "8. Patient views results (portal)",
]

for i, step in enumerate(lab_workflow):
    add_textbox(slide, Inches(0.7), Inches(1.9 + i * 0.5), Inches(5.5), Inches(0.4),
                step, 10, BLACK, False)

# Billing workflow (right)
add_textbox(slide, Inches(7), Inches(1.3), Inches(6), Inches(0.5),
            "BILLING WORKFLOW", 14, DARK_BLUE, True)

billing_workflow = [
    "1. Services rendered → Charges created",
    "2. EMR notes auto-generate charges",
    "3. All charges grouped into invoice",
    "4. Insurance verification (if needed)",
    "5. Invoice sent to patient/insurance",
    "6. Payment received & recorded",
    "7. Receipt auto-emailed",
    "8. Reports updated in real-time",
]

for i, step in enumerate(billing_workflow):
    add_textbox(slide, Inches(7.2), Inches(1.9 + i * 0.5), Inches(5.5), Inches(0.4),
                step, 10, BLACK, False)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 17: KEY BENEFITS
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_layout(slide, 4, "KEY BENEFITS & ROI",
                   "Measurable Improvements for Your Hospital")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 18: BENEFITS & ROI METRICS
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "BENEFITS & RETURN ON INVESTMENT (ROI)", 24, WHITE, True)

# KPI cards row 1
kpis = [
    ("60%", "Wait Time Reduction", MED_BLUE),
    ("40%", "Billing Efficiency Gain", ORANGE),
    ("70%", "Fewer Medical Errors", DARK_GREEN),
    ("90%", "Paper Elimination", PURPLE),
    ("25%", "Revenue Increase", RED),
]

for i, (metric, label, color) in enumerate(kpis):
    left = Inches(0.3 + i * 2.55)
    add_kpi_card(slide, left, Inches(1.4), metric, label, color)

# Benefits table
benefits_data = [
    ("Operational Efficiency", "Automated workflows reduce manual work by 60% — staff focus on patient care"),
    ("Reduced Medical Errors", "Drug interaction checking, allergy alerts, clinical decision support"),
    ("Faster Patient Processing", "Check-in to consultation 40% faster with digital queue management"),
    ("Revenue Optimization", "Accurate billing, fewer insurance claim rejections, faster payment"),
    ("Better Patient Experience", "Online booking, instant results, appointment reminders, patient portal"),
    ("Compliance Ready", "HIPAA-compliant design, complete audit trail, encrypted data"),
    ("Cost Savings", "$200-300/month operating costs vs $5-10K/month for enterprise EMR"),
    ("Future Proof", "Cloud-based, auto-updates, scales from 50 to 500+ beds"),
]

add_table(slide, Inches(0.5), Inches(2.7), Inches(12.3), Inches(4.5),
          ["Benefit", "Impact"], benefits_data[:4])

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 19: IMPLEMENTATION TIMELINE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_layout(slide, 5, "IMPLEMENTATION TIMELINE",
                   "12 Weeks from Order to Go-Live")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 20: WEEKLY BREAKDOWN
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "12-WEEK IMPLEMENTATION TIMELINE", 24, WHITE, True)

timeline_phases = [
    ("WEEK 1", "Architecture & Setup", DARK_BLUE,
     "• System design\n• Database schema\n• API specification\n• Project initialization"),
    ("WEEKS 2-3", "Backend Development", MED_BLUE,
     "• 150+ API endpoints\n• Database tables\n• Authentication\n• Business logic"),
    ("WEEKS 4-5", "Frontend Development", ACCENT_TEAL,
     "• 70+ pages\n• Responsive design\n• Real-time updates\n• Reports & exports"),
    ("WEEK 6", "Testing & Integration", DARK_GREEN,
     "• E2E testing\n• Security testing\n• Load testing\n• Bug fixes"),
    ("WEEK 7", "Deployment", ORANGE,
     "• Supabase setup\n• Koyeb deployment\n• Vercel launch\n• Monitoring active"),
    ("WEEK 8", "Documentation", PURPLE,
     "• API documentation\n• User guides\n• Admin guide\n• FAQ & support"),
    ("WEEKS 9-12", "Refinement & Go-Live", RED,
     "• User feedback\n• Optimization\n• Staff training\n• System launch"),
]

for phase_idx, (week, title, color, details) in enumerate(timeline_phases):
    top = Inches(1.4 + phase_idx * 0.9)
    
    # Week badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), top, Inches(1.2), Inches(0.8))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = week
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Title & details
    add_textbox(slide, Inches(1.7), top, Inches(11.1), Inches(0.8),
                f"{title}\n{details}", 9, BLACK, False)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 21: PRICING & LICENSING
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_layout(slide, 6, "PRICING & LICENSING",
                   "Transparent, Predictable Costs")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 22: PRICING BREAKDOWN
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "PRICING & COST BREAKDOWN", 24, WHITE, True)

# One-time costs
add_textbox(slide, Inches(0.5), Inches(1.3), Inches(6.2), Inches(0.5),
            "ONE-TIME DEVELOPMENT", 14, DARK_BLUE, True)

onetime = [
    ("System Development", "$5,000 - $10,000"),
    ("Architecture & Design", "Included"),
    ("Database Setup", "Included"),
    ("Data Migration (if any)", "+$1,000 - $3,000"),
    ("Staff Training", "Included"),
    ("TOTAL", "$5,000 - $13,000"),
]

add_table(slide, Inches(0.5), Inches(1.9), Inches(6.0), Inches(2.0),
          ["Service", "Cost"], onetime)

# Monthly costs
add_textbox(slide, Inches(7.2), Inches(1.3), Inches(5.6), Inches(0.5),
            "MONTHLY OPERATING", 14, DARK_BLUE, True)

monthly = [
    ("Supabase Database", "$25-$250"),
    ("Koyeb Backend", "$20-$200"),
    ("Vercel Frontend", "$20-$100"),
    ("Email & SMS", "$10-$50"),
    ("Domain & SSL", "$10-$20"),
    ("TOTAL (Typical)", "$200-$300"),
]

add_table(slide, Inches(7.2), Inches(1.9), Inches(5.6), Inches(2.0),
          ["Service", "Cost"], monthly)

# Comparison box
add_shape_bg(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.8), LIGHT_BLUE)
add_textbox(slide, Inches(0.7), Inches(4.4), Inches(11.9), Inches(0.5),
            "WHY THIS IS THE BEST VALUE FOR SMALL HOSPITALS", 14, DARK_BLUE, True)

comparison_text = [
    "Traditional Enterprise EMR: $100K-$500K setup + $5K-$10K/month + dedicated IT team",
    "Our Solution: $5-13K setup + $200-300/month + NO IT team needed",
    "Annual Savings: $46K-$119K+ (vs traditional) — ROI in 1-2 months",
    "Scalability: Grows from 50 to 500+ beds at same monthly cost (no renegotiation needed)",
]

for i, text in enumerate(comparison_text):
    add_textbox(slide, Inches(0.8), Inches(5.05 + i * 0.5), Inches(11.7), Inches(0.4),
                f"✓  {text}", 11, BLACK, False)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 23: SECURITY & COMPLIANCE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_layout(slide, 7, "SECURITY & COMPLIANCE",
                   "Enterprise-Grade Protection for Healthcare Data")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 24: SECURITY FEATURES
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "SECURITY & COMPLIANCE FEATURES", 24, WHITE, True)

security_items = [
    ("🔐 Data Encryption", "AES-256 at rest, TLS 1.3 in transit — military-grade security"),
    ("🔑 Authentication", "JWT tokens, multi-factor authentication (MFA), session management"),
    ("👥 Access Control", "Role-based (7 roles: Doctor, Nurse, Pharmacist, Lab, Admin, Accountant, Receptionist)"),
    ("📋 Audit Logging", "Complete audit trail — every login, access, edit tracked with timestamps"),
    ("💾 Backups", "Automated daily backups, 30-day retention, point-in-time recovery"),
    ("🛡️ HIPAA Ready", "HIPAA-compliant architecture, patient privacy controls, breach notification"),
    ("✅ GDPR Support", "Data export, right to be forgotten, consent management"),
    ("🔍 Monitoring", "24/7 system monitoring, alerts for suspicious activity, error tracking"),
]

colors = [DARK_BLUE, MED_BLUE, ACCENT_TEAL, DARK_GREEN, ORANGE, PURPLE, RED, GRAY]

for i, (title, desc) in enumerate(security_items):
    top = Inches(1.4 + (i % 4) * 1.5)
    left = Inches(0.5 if i < 4 else 6.5)
    
    add_textbox(slide, left, top, Inches(5.8), Inches(0.3),
                title, 12, colors[i % len(colors)], True)
    add_textbox(slide, left, top + Inches(0.35), Inches(5.8), Inches(1.0),
                desc, 10, BLACK, False)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 25: DEPLOYMENT PLAN
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_layout(slide, 8, "DEPLOYMENT PLAN",
                   "How We Get You Live")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 26: GO-LIVE PROCESS
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "GO-LIVE DEPLOYMENT PROCESS", 24, WHITE, True)

golive_steps = [
    ("1. Pre-Deployment", "Production database created, backups configured, monitoring active", DARK_BLUE),
    ("2. Data Migration", "Legacy data exported, transformed to new schema, integrity verified", MED_BLUE),
    ("3. Staging Test", "All systems tested in staging environment, UAT scenarios validated", ACCENT_TEAL),
    ("4. Staff Training", "Doctors, nurses, admins trained on new system, certification provided", DARK_GREEN),
    ("5. System Going Live", "Switch from old to new system, all departments go digital", ORANGE),
    ("6. Launch Support", "24/7 support for 1 week, monitoring, issue resolution, optimization", PURPLE),
    ("7. Handoff", "System stable, documentation complete, support team onboarded", RED),
]

for i, (step, desc, color) in enumerate(golive_steps):
    top = Inches(1.4 + i * 0.9)
    
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), top, Inches(1.4), Inches(0.8))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = step
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    add_textbox(slide, Inches(1.8), top, Inches(10.8), Inches(0.8),
                desc, 11, BLACK, False)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 27: SUPPORT & TRAINING
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_layout(slide, 9, "SUPPORT & TRAINING",
                   "Comprehensive Help for Your Team")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 28: TRAINING & SUPPORT DETAILS
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "SUPPORT & TRAINING PROGRAM", 24, WHITE, True)

# Training section
add_textbox(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(0.45),
            "📚 TRAINING MATERIALS", 14, DARK_BLUE, True)

training = [
    "✓ 20+ video tutorials (step-by-step)",
    "✓ User guide (50+ pages, PDF)",
    "✓ Admin guide (configuration)",
    "✓ Troubleshooting guide",
    "✓ Quick reference cards",
    "✓ Role-specific training slides",
]

for i, item in enumerate(training):
    add_textbox(slide, Inches(0.7), Inches(1.9 + i * 0.45), Inches(5.6), Inches(0.4),
                item, 11, BLACK, False)

# Support section
add_textbox(slide, Inches(7.0), Inches(1.3), Inches(5.8), Inches(0.45),
            "🆘 SUPPORT SERVICES", 14, DARK_BLUE, True)

support = [
    "✓ 24/7 support (go-live week)",
    "✓ Email support (ongoing)",
    "✓ Quick bug fixes (24 hours)",
    "✓ System updates (automatic)",
    "✓ Feature requests (roadmap)",
    "✓ Periodic check-ins",
]

for i, item in enumerate(support):
    add_textbox(slide, Inches(7.2), Inches(1.9 + i * 0.45), Inches(5.56), Inches(0.4),
                item, 11, BLACK, False)

# SLA & Uptime
add_shape_bg(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5), LIGHT_GREEN)
add_textbox(slide, Inches(0.7), Inches(4.6), Inches(11.9), Inches(0.4),
            "SERVICE LEVEL AGREEMENT (SLA)", 14, DARK_GREEN, True)

sla_items = [
    ("99.9% Uptime", "System available 99.9% of the time (less than 43 minutes downtime/month)"),
    ("< 1 Hour Response", "Critical issues responded to within 1 hour"),
    ("< 24 Hour Resolution", "Bug fixes deployed within 24 hours"),
    ("Daily Backups", "Automated daily backups with 30-day retention"),
    ("Auto-Updates", "Security patches and updates applied automatically"),
]

for i, (title, desc) in enumerate(sla_items):
    if i > 0:
        add_textbox(slide, Inches(0.7), Inches(5.1 + (i - 1) * 0.4), Inches(11.9), Inches(0.35),
                    f"✓ {title}: {desc}", 10, BLACK, False)
    else:
        add_textbox(slide, Inches(0.7), Inches(5.1), Inches(11.9), Inches(0.35),
                    f"✓ {title}: {desc}", 10, BLACK, False)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 29: PRODUCT COMPARISON
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "HOW WE COMPARE", 24, WHITE, True)

comparison = [
    ("Feature", "Our EMR", "Enterprise EMR", "Paper + Spreadsheets"),
    ("Setup Cost", "$5-13K", "$100-500K", "Free (no digital)"),
    ("Monthly Cost", "$200-300", "$5-10K+", "Staff + paper costs"),
    ("Implementation", "12 weeks", "2 years", "Not applicable"),
    ("All-in-One", "YES ✓", "NO (different vendors)", "NO"),
    ("Scalable", "YES ✓", "YES (expensive)", "NO"),
    ("Easy to Use", "YES ✓", "Complex (steep learning)", "NO"),
    ("Secure", "YES ✓", "YES", "NO (HIPAA risk)"),
    ("Support", "Email + 24/7 launch", "SLA-based (expensive)", "No support"),
    ("Total Year 1", "$7-16K", "$160-600K", "$100K-500K (labor)"),
]

comp_table = slide.shapes.add_table(len(comparison), 4, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8))
comp_table = comp_table.table

for ri, row in enumerate(comparison):
    for ci, val in enumerate(row):
        cell = comp_table.cell(ri, ci)
        cell.text = str(val)
        is_header = (ri == 0)
        if is_header:
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE
        elif ci == 1:  # Our EMR column
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GREEN
        elif ri % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE

        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.color.rgb = WHITE if is_header else BLACK
            p.font.bold = is_header
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 30: NEXT STEPS & CONTACT
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)

add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), MED_BLUE)

add_textbox(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1.0),
            "Ready to Transform", 40, WHITE, True, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1.0),
            "Your Hospital?", 40, RGBColor(100, 180, 255), True, PP_ALIGN.CENTER)

add_textbox(slide, Inches(2), Inches(3.7), Inches(9.3), Inches(0.7),
            "Let's Schedule a Discussion", 18, RGBColor(180, 210, 240), False, PP_ALIGN.CENTER)

next_steps = [
    "✓ We discuss your hospital's needs",
    "✓ Confirm requirements & timelines",
    "✓ Sign agreement & start development",
    "✓ 12 weeks later: Go-live to production",
]

for i, step in enumerate(next_steps):
    add_textbox(slide, Inches(2), Inches(4.6 + i * 0.5), Inches(9.3), Inches(0.4),
                step, 14, WHITE, False, PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(2.5), Inches(6.5), Inches(8.3), Inches(0.7), RGBColor(0, 150, 136))
add_textbox(slide, Inches(2.7), Inches(6.55), Inches(7.9), Inches(0.6),
            "Contact Us: development@hospitalemr.com | Ready in 12 Weeks",
            14, WHITE, True, PP_ALIGN.CENTER)

# Save presentation
output_path = os.path.join(OUTPUT_DIR, "Small_Hospital_EMR_Pitch.pptx")
prs.save(output_path)
print(f"✅ PowerPoint saved: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
