"""
Generate comprehensive Hospital EMR Pitch PowerPoint Presentation.
50+ slides with professional styling, diagrams, tables, and flowcharts.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "Hospital_EMR_Pitch")
os.makedirs(OUTPUT_DIR, exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Color Palette ────────────────────────────────────────────────────
DARK_BLUE = RGBColor(0, 51, 102)
MED_BLUE = RGBColor(0, 102, 153)
LIGHT_BLUE = RGBColor(230, 240, 250)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(33, 33, 33)
DARK_GREEN = RGBColor(0, 102, 51)
LIGHT_GREEN = RGBColor(230, 245, 230)
RED = RGBColor(180, 30, 30)
ORANGE = RGBColor(230, 115, 0)
PURPLE = RGBColor(102, 0, 102)
GRAY = RGBColor(100, 100, 100)
LIGHT_GRAY = RGBColor(240, 240, 240)
MED_GRAY = RGBColor(200, 200, 200)
ACCENT_TEAL = RGBColor(0, 150, 136)

# ─── Helper Functions ─────────────────────────────────────────────────

def add_bg(slide, color=DARK_BLUE):
    """Set solid background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=BLACK, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_multi_text(slide, left, top, width, height, lines,
                   font_size=14, color=BLACK, spacing=1.2, bullet=False):
    """Add multiple lines of text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line, tuple):
            text, line_bold, line_color = line
        else:
            text = line
            line_bold = False
            line_color = color
        if bullet and not text.startswith("•"):
            text = "• " + text
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = line_color
        p.font.bold = line_bold
        p.font.name = "Calibri"
        p.space_after = Pt(font_size * spacing * 0.4)
    return txBox


def add_table(slide, left, top, width, height, headers, rows):
    """Add a formatted table."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    # Column widths (distribute evenly if not specified)
    col_w = int(width / num_cols)
    for i in range(num_cols):
        table.columns[i].width = col_w

    # Headers
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            if r_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BLUE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = BLACK
                p.font.name = "Calibri"
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


def add_flowchart_boxes(slide, steps, start_left, start_top, box_w, box_h, gap, color=DARK_BLUE):
    """Add horizontal or vertical flowchart boxes with arrows."""
    for i, step in enumerate(steps):
        left = start_left
        top = start_top + i * (box_h + gap)
        # Box
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, box_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(10)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(2)

        # Arrow (except last)
        if i < len(steps) - 1:
            arrow_top = top + box_h
            arrow_shape = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW, left + box_w // 2 - Inches(0.15),
                arrow_top, Inches(0.3), gap
            )
            arrow_shape.fill.solid()
            arrow_shape.fill.fore_color.rgb = MED_BLUE
            arrow_shape.line.fill.background()


def add_horizontal_flow(slide, steps, start_left, start_top, box_w, box_h, gap, color=DARK_BLUE):
    """Add horizontal flow boxes."""
    for i, step in enumerate(steps):
        left = start_left + i * (box_w + gap)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, start_top, box_w, box_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(9)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                left + box_w, start_top + box_h // 2 - Inches(0.12),
                gap, Inches(0.24)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MED_BLUE
            arrow.line.fill.background()


def title_slide_layout(slide, section_num, title, subtitle=""):
    """Create a section title slide with dark blue background."""
    add_bg(slide, DARK_BLUE)
    # Section number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(5.8), Inches(1.5), Inches(1.6), Inches(1.6)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = WHITE
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = str(section_num)
    p.font.size = Pt(48)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(12)

    add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(1.5),
                title, 40, WHITE, True, PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide, Inches(1), Inches(5), Inches(11), Inches(1),
                    subtitle, 20, RGBColor(180, 210, 240), False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1: COVER
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BLUE)

# Top accent bar
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), MED_BLUE)

# Hospital icon box
icon_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(1.0), Inches(2.3), Inches(1.5)
)
icon_box.fill.solid()
icon_box.fill.fore_color.rgb = WHITE
icon_box.line.fill.background()
tf = icon_box.text_frame
p = tf.paragraphs[0]
p.text = "🏥"
p.font.size = Pt(60)
p.alignment = PP_ALIGN.CENTER

# Title
add_textbox(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(1.2),
            "HOSPITAL EMR SYSTEM", 44, WHITE, True, PP_ALIGN.CENTER)

# Subtitle
add_textbox(slide, Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.8),
            "Electronic Medical Records & Hospital Management Platform", 22,
            RGBColor(180, 210, 240), False, PP_ALIGN.CENTER)

# Tagline
add_textbox(slide, Inches(2), Inches(5.2), Inches(9.3), Inches(0.7),
            "Comprehensive Solution for Private Hospital Operations", 18,
            RGBColor(150, 180, 210), False, PP_ALIGN.CENTER, "Calibri")

# Bottom info bar
add_shape_bg(slide, Inches(0), Inches(6.6), prs.slide_width, Inches(0.9), RGBColor(0, 40, 80))
add_textbox(slide, Inches(1), Inches(6.7), Inches(11.3), Inches(0.6),
            "CONFIDENTIAL PITCH DOCUMENT   |   February 2026   |   Version 1.0",
            14, RGBColor(180, 200, 220), False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2: TABLE OF CONTENTS
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "TABLE OF CONTENTS", 32, WHITE, True, PP_ALIGN.LEFT)

toc_left = [
    "1.  Executive Summary",
    "2.  System Overview & Architecture",
    "3.  Core Modules (15+ Modules)",
    "4.  Administrative Modules",
    "5.  Workflow Diagrams",
    "6.  Integration & Interoperability",
]
toc_right = [
    "7.   Security & Compliance",
    "8.   Technology Stack",
    "9.   Implementation Plan",
    "10.  Pricing & Licensing",
    "11.  Support & Maintenance",
    "12.  Why Choose Our EMR?",
]

add_multi_text(slide, Inches(1), Inches(1.5), Inches(5.5), Inches(5),
               toc_left, 18, DARK_BLUE, 2.0)
add_multi_text(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(5),
               toc_right, 18, DARK_BLUE, 2.0)

# Sub-modules list
sub_modules = "Patient Mgmt • EMR/EHR • Appointments • OPD • IPD/Wards • Emergency • Operating Theatre • Nursing • Laboratory • Pharmacy • Radiology • Blood Bank • Inventory • Billing • HR/Staff"
add_textbox(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.8),
            sub_modules, 11, GRAY, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3: EXECUTIVE SUMMARY - SECTION TITLE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_layout(slide, 1, "EXECUTIVE SUMMARY", "The Big Picture — Why This EMR System?")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4: EXECUTIVE SUMMARY - CONTENT
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "1. EXECUTIVE SUMMARY", 28, WHITE, True)

add_textbox(slide, Inches(0.5), Inches(1.3), Inches(12), Inches(1.0),
            "Our Hospital EMR system is a comprehensive, integrated healthcare management platform "
            "designed specifically for private hospitals. It digitizes every aspect of hospital "
            "operations — from patient registration to discharge, from pharmacy dispensing to "
            "laboratory results, from staff scheduling to financial reporting.",
            13, BLACK, False, PP_ALIGN.LEFT)

# Key Value Propositions - 2x4 grid of boxes
kvp = [
    ("📄 Paperless Operations", "Eliminate paper records entirely"),
    ("🛡️ Patient Safety", "Drug interaction & allergy alerts"),
    ("💰 Revenue Optimization", "Automated billing & claims"),
    ("⚡ Operational Efficiency", "60% wait time reduction"),
    ("✅ Regulatory Compliance", "Full audit trail & encryption"),
    ("📊 Data-Driven Decisions", "Real-time dashboards & BI"),
    ("👥 Staff Productivity", "Automated scheduling & tasks"),
    ("😊 Patient Satisfaction", "Portal, booking, notifications"),
]

for i, (title, desc) in enumerate(kvp):
    col = i % 4
    row = i // 4
    left = Inches(0.5 + col * 3.1)
    top = Inches(2.7 + row * 1.9)
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.9), Inches(1.6)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BLUE
    box.line.color.rgb = MED_BLUE
    box.line.width = Pt(1)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(10)
    p2.font.color.rgb = GRAY
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5: ROI HIGHLIGHTS
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "RETURN ON INVESTMENT (ROI)", 28, WHITE, True)

roi_data = [
    ("60%", "Wait Time\nReduction", MED_BLUE),
    ("90%", "Fewer Medical\nErrors", DARK_GREEN),
    ("25%", "Revenue\nIncrease", ORANGE),
    ("30%", "Staff Efficiency\nImprovement", ACCENT_TEAL),
    ("40%", "Faster Insurance\nClaims", PURPLE),
    ("$200K+", "Annual Paper\nCost Savings", RED),
]

for i, (metric, label, color) in enumerate(roi_data):
    left = Inches(0.4 + i * 2.1)
    top = Inches(1.5)
    # Circle for metric
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left + Inches(0.25), top, Inches(1.5), Inches(1.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = metric
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(18)

    add_textbox(slide, left, top + Inches(1.7), Inches(2.0), Inches(0.8),
                label, 12, BLACK, True, PP_ALIGN.CENTER)

# Bottom comparison
add_shape_bg(slide, Inches(0.5), Inches(3.8), Inches(12.3), Inches(3.2), LIGHT_GRAY)
add_textbox(slide, Inches(0.7), Inches(3.9), Inches(12), Inches(0.5),
            "Before vs After EMR Implementation", 16, DARK_BLUE, True, PP_ALIGN.CENTER)

before_items = [
    "❌ Paper-based records — lost, damaged, misfiled",
    "❌ Manual billing — errors, delays, revenue leakage",
    "❌ No real-time data — decisions based on guesswork",
    "❌ Long patient wait times — poor satisfaction",
    "❌ Siloed departments — no information sharing",
]
after_items = [
    "✅ Digital records — instant access, never lost",
    "✅ Automated billing — accurate, fast, complete capture",
    "✅ Real-time dashboards — data-driven decisions",
    "✅ Streamlined workflows — 60% shorter wait times",
    "✅ Integrated system — all departments connected",
]

add_multi_text(slide, Inches(0.8), Inches(4.4), Inches(5.8), Inches(2.5),
               before_items, 11, BLACK, 1.5)
add_multi_text(slide, Inches(6.8), Inches(4.4), Inches(5.8), Inches(2.5),
               after_items, 11, DARK_GREEN, 1.5)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6: SYSTEM ARCHITECTURE - SECTION TITLE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_layout(slide, 2, "SYSTEM OVERVIEW & ARCHITECTURE",
                   "Modern Three-Tier Architecture")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7: ARCHITECTURE DIAGRAM
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "HIGH-LEVEL SYSTEM ARCHITECTURE", 28, WHITE, True)

# Client layer
clients = ["🖥️ Web Browser\n(Desktop)", "📱 Mobile App\n(iOS/Android)", "📟 Tablet/Kiosk\n(Check-in)"]
for i, cl in enumerate(clients):
    left = Inches(2.8 + i * 2.8)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.4), Inches(2.2), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = MED_BLUE
    box.line.color.rgb = DARK_BLUE
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = cl
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

# API Gateway
gw = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.8), Inches(2.7), Inches(5.7), Inches(0.7))
gw.fill.solid()
gw.fill.fore_color.rgb = ORANGE
gw.line.fill.background()
tf = gw.text_frame
p = tf.paragraphs[0]
p.text = "API GATEWAY / LOAD BALANCER  (NGINX + SSL + Rate Limiting)"
p.font.size = Pt(12)
p.font.color.rgb = WHITE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Services
services = [
    ("🔐 Auth Service\nJWT / RBAC / MFA", DARK_GREEN),
    ("⚕️ Core EMR API\nREST + GraphQL", DARK_BLUE),
    ("📊 Analytics Service\nReports + BI", PURPLE),
]
for i, (svc, color) in enumerate(services):
    left = Inches(2.5 + i * 3.0)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3.8), Inches(2.6), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = svc
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

# Database layer
dbs = [
    ("🗄️ PostgreSQL\nPrimary Database", DARK_BLUE),
    ("⚡ Redis\nCache + Sessions", RED),
    ("🔍 Elasticsearch\nFull-Text Search", ACCENT_TEAL),
    ("📁 MinIO / S3\nFile Storage", ORANGE),
    ("🏥 Orthanc\nDICOM/PACS", PURPLE),
]
for i, (db, color) in enumerate(dbs):
    left = Inches(1.0 + i * 2.4)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(5.2), Inches(2.1), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = db
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

# Layer labels
add_textbox(slide, Inches(0.3), Inches(1.5), Inches(2.3), Inches(0.5),
            "PRESENTATION", 11, MED_BLUE, True)
add_textbox(slide, Inches(0.3), Inches(2.8), Inches(2.3), Inches(0.5),
            "GATEWAY", 11, ORANGE, True)
add_textbox(slide, Inches(0.3), Inches(3.9), Inches(2.3), Inches(0.5),
            "APPLICATION", 11, DARK_BLUE, True)
add_textbox(slide, Inches(0.3), Inches(5.3), Inches(2.3), Inches(0.5),
            "DATA LAYER", 11, DARK_GREEN, True)

# Arrows between layers (simple rectangles as connectors)
for top_pos in [Inches(2.35), Inches(3.45), Inches(4.75)]:
    arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.5), top_pos, Inches(0.3), Inches(0.3))
    arr.fill.solid()
    arr.fill.fore_color.rgb = MED_GRAY
    arr.line.fill.background()

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8: SYSTEM COMPONENTS TABLE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "SYSTEM COMPONENTS", 28, WHITE, True)

components = [
    ("Frontend (Web)", "React.js / Next.js 14", "SSR, responsive, PWA support"),
    ("Mobile App", "React Native / Flutter", "Cross-platform iOS & Android, offline"),
    ("Backend API", "Node.js (NestJS) / Python (FastAPI)", "REST + GraphQL + WebSocket"),
    ("Database", "PostgreSQL 16+", "ACID, JSON support, full-text search"),
    ("Caching", "Redis 7+", "Sessions, real-time data, pub/sub"),
    ("Search Engine", "Elasticsearch 8+", "Patient & clinical term search"),
    ("File Storage", "MinIO / AWS S3", "Medical images, documents, DICOM"),
    ("Message Queue", "RabbitMQ / Redis Streams", "Async processing, notifications"),
    ("PACS Server", "Orthanc / dcm4chee", "DICOM medical imaging"),
    ("CI/CD", "GitHub Actions / GitLab CI", "Automated testing & deployment"),
    ("Monitoring", "Prometheus + Grafana", "System health, performance metrics"),
    ("Logging", "ELK Stack", "Centralized log management"),
]
add_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.5),
          ["Component", "Technology", "Purpose"], components)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9: CORE MODULES - SECTION TITLE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_layout(slide, 3, "CORE MODULES & FEATURES",
                   "15+ Integrated Modules Covering Every Department")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10: MODULE OVERVIEW MAP
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "MODULE OVERVIEW MAP", 28, WHITE, True)

# Module categories in colored boxes
categories = [
    ("CLINICAL MODULES", DARK_BLUE, [
        "Patient Management", "EMR / EHR", "OPD Management",
        "IPD / Wards", "Emergency Dept", "Operating Theatre", "Nursing Station"
    ]),
    ("SUPPORT MODULES", DARK_GREEN, [
        "Laboratory (LIS)", "Pharmacy", "Radiology (RIS/PACS)",
        "Blood Bank", "Inventory & Supply", "Diet / Kitchen"
    ]),
    ("FINANCIAL MODULES", ORANGE, [
        "Billing & Invoicing", "Insurance Claims",
        "Revenue Cycle", "Accounting Integration"
    ]),
    ("ADMIN & HR MODULES", PURPLE, [
        "Admin Dashboard", "Reports & BI", "User Management",
        "Staff / HR", "Audit Trail", "Settings"
    ]),
]

for cat_idx, (cat_name, cat_color, modules) in enumerate(categories):
    col = cat_idx % 2
    row = cat_idx // 2
    left = Inches(0.5 + col * 6.4)
    top = Inches(1.3 + row * 3.0)

    # Category header
    hdr = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(6.0), Inches(0.5)
    )
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = cat_color
    hdr.line.fill.background()
    tf = hdr.text_frame
    p = tf.paragraphs[0]
    p.text = cat_name
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Module items
    for m_idx, mod in enumerate(modules):
        m_col = m_idx % 3
        m_row = m_idx // 3
        m_left = left + Inches(m_col * 2.0)
        m_top = top + Inches(0.6 + m_row * 0.65)
        m_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, m_left, m_top, Inches(1.9), Inches(0.5)
        )
        m_box.fill.solid()
        m_box.fill.fore_color.rgb = LIGHT_BLUE if cat_idx < 2 else LIGHT_GREEN
        m_box.line.color.rgb = cat_color
        m_box.line.width = Pt(1)
        tf = m_box.text_frame
        p = tf.paragraphs[0]
        p.text = mod
        p.font.size = Pt(8)
        p.font.color.rgb = cat_color
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════════════════════════════════════════
# SLIDES 11-25: INDIVIDUAL MODULE DETAIL SLIDES
# ═══════════════════════════════════════════════════════════════════════

module_slides = [
    ("3.1 PATIENT MANAGEMENT", [
        ("Patient Registration", "Demographics, contacts, emergency contacts, photo, biometric ID"),
        ("Unique Patient ID (UHID)", "Auto-generated with barcode/QR code for all encounters"),
        ("Patient Search", "Instant search by name, ID, phone, national ID"),
        ("Patient Categories", "VIP, regular, insurance, corporate, pediatric, geriatric"),
        ("Medical History", "Past illnesses, surgeries, family history, allergies"),
        ("Insurance Management", "Multiple policies, pre-authorization, coverage limits"),
        ("Document Upload", "ID cards, insurance, referral letters, consent forms"),
        ("Patient Portal", "Self-service: view records, book appointments, download reports"),
        ("Communication", "SMS, email, WhatsApp for appointments & results"),
    ]),
    ("3.2 ELECTRONIC MEDICAL RECORDS (EMR)", [
        ("Clinical Notes (SOAP)", "Subjective, Objective, Assessment, Plan — with templates"),
        ("Diagnosis (ICD-10/11)", "Searchable ICD coding, multiple diagnoses per encounter"),
        ("E-Prescribing", "Drug interaction checks, dosage calculators, generic substitution"),
        ("Vitals Recording", "BP, temperature, pulse, SpO2, weight, height, BMI auto-calc"),
        ("Order Management", "Lab orders, radiology orders, procedure orders — all digital"),
        ("Clinical Templates", "Specialty-specific: cardiology, orthopedics, OB/GYN, peds"),
        ("Allergy & Alert System", "Drug allergies, food allergies, condition alerts"),
        ("Medical Timeline", "Chronological view of all encounters, results, procedures"),
        ("Clinical Decision Support", "Evidence-based alerts, risk scores, guidelines"),
        ("Voice-to-Text / E-Signatures", "AI dictation + digital signatures on documents"),
    ]),
    ("3.3 APPOINTMENT & SCHEDULING", [
        ("Online Booking", "Patients book via web portal or mobile app — 24/7"),
        ("Walk-in Queue", "Real-time queue management with estimated wait times"),
        ("Doctor Schedule", "Define availability, breaks, leave, on-call schedules"),
        ("Multi-Resource Booking", "Simultaneously book doctor + room + equipment"),
        ("Recurring Appointments", "Auto-schedule for dialysis, physiotherapy, etc."),
        ("Reminders", "SMS/email/WhatsApp reminders 24h and 1h before"),
        ("Telemedicine Slots", "Video consultation with integrated calling"),
        ("Capacity Dashboard", "Real-time clinic capacity and utilization view"),
    ]),
    ("3.4 STAFF & HR MANAGEMENT", [
        ("Employee Profiles", "Credentials, certifications, specializations"),
        ("Attendance & Time", "Clock-in/out, biometric, overtime tracking"),
        ("Shift Scheduling", "Auto-rosters, swap requests, on-call management"),
        ("Leave Management", "Apply, approve, track — annual, sick, maternity, emergency"),
        ("Payroll Integration", "Salary calculation, deductions, bonuses, pay slips"),
        ("Performance Reviews", "360° evaluations, KPIs, competency assessments"),
        ("Training & CME", "Continuing education, certifications, expiry alerts"),
        ("Credentialing", "Medical licenses, board certifications, verification"),
    ]),
    ("3.5 BILLING & FINANCIAL MANAGEMENT", [
        ("Patient Billing", "Auto-generate from orders, procedures, bed charges"),
        ("Insurance Claims", "Electronic submission, pre-auth, adjudication tracking"),
        ("Multiple Payments", "Cash, card, mobile money, bank transfer, installments"),
        ("Package Billing", "Health check, surgery, maternity packages"),
        ("Revenue Dashboard", "Daily/weekly/monthly revenue, department collections"),
        ("Tax Management", "VAT/GST, tax invoices, exemptions"),
        ("Financial Reports", "P&L, balance sheet, aging, revenue forecasts"),
        ("Accounting Integration", "QuickBooks, Xero, Tally, SAP export"),
    ]),
    ("3.6 PHARMACY MANAGEMENT", [
        ("E-Prescription Receipt", "Receive prescriptions electronically from EMR"),
        ("Drug Dispensing", "Barcode-based, patient verification, label printing"),
        ("Drug Interaction Check", "Real-time drug-drug & drug-allergy alerts"),
        ("Inventory Management", "Stock levels, reorder points, expiry, batch tracking"),
        ("Controlled Substances", "Special tracking for narcotics and controlled meds"),
        ("Formulary", "Drug formulary with generic substitution suggestions"),
        ("Purchase Orders", "Auto-PO generation when stock hits reorder level"),
        ("Point of Sale", "Walk-in sales, OTC medications, billing"),
    ]),
    ("3.7 LABORATORY MANAGEMENT (LIS)", [
        ("Order Management", "Receive digital lab orders, priority flagging (STAT)"),
        ("Sample Collection", "Barcode labels, phlebotomy scheduling, tracking"),
        ("Sample Tracking", "Full chain: collection → processing → reporting"),
        ("Result Entry", "Manual or instrument-interfaced with normal ranges"),
        ("Auto-Validation", "Rules-based validation for normal range results"),
        ("Critical Alerts", "Immediate notification for critical/panic values"),
        ("Quality Control", "QC tracking, Levey-Jennings, Westgard rules"),
        ("Instrument Interface", "Bi-directional connection with analyzers"),
        ("TAT Monitoring", "Turnaround time tracking, bottleneck identification"),
    ]),
    ("3.8 RADIOLOGY & IMAGING (RIS/PACS)", [
        ("Order Management", "Digital radiology orders with clinical indication"),
        ("DICOM Integration", "Modality Worklist, image storage, PACS server"),
        ("Web-Based Viewer", "DICOM viewer with measurement & comparison tools"),
        ("Structured Reports", "Templates, voice dictation, automated distribution"),
        ("Equipment Tracking", "Maintenance schedules, utilization, downtime reports"),
    ]),
    ("3.9 INVENTORY & SUPPLY CHAIN", [
        ("Multi-Store Mgmt", "Central store, department sub-stores, pharmacy, kitchen"),
        ("Purchase Workflow", "Requisition → PO → GRN → Issue with approvals"),
        ("Expiry Management", "FEFO/FIFO tracking, alerts at 30/60/90 days"),
        ("ABC/VED Analysis", "Classify items by value and criticality"),
        ("Barcode/RFID", "Item identification and tracking throughout hospital"),
        ("Vendor Management", "Evaluation, price comparison, contract management"),
    ]),
    ("3.10 NURSING STATION", [
        ("Patient Worklist", "Assigned patients with alerts and pending tasks"),
        ("Vital Signs Charting", "Record/chart vitals with auto-alerts for abnormals"),
        ("eMAR", "Barcode scan patient + med for safe medication administration"),
        ("Care Plans", "Create/manage individualized nursing care plans"),
        ("Intake/Output", "Track fluid intake, urine output, drain output"),
        ("Shift Handover", "Structured handover notes with critical information"),
    ]),
    ("3.11 IN-PATIENT (WARD) MANAGEMENT", [
        ("Admission", "Direct, emergency, or transfer admission with bed assignment"),
        ("Bed Management", "Visual bed map, occupancy dashboard, ICU/general/private"),
        ("Ward Dashboard", "Real-time ward status, patient list, alerts"),
        ("Transfer Management", "Ward-to-ward, bed-to-bed with charge adjustments"),
        ("Discharge Planning", "Summary, follow-up, take-home meds, clearance checklist"),
        ("Census Reports", "Daily census, ADT reports, length of stay tracking"),
    ]),
    ("3.12 OPD & 3.13 EMERGENCY DEPT", [
        ("OPD: Check-in/Queue", "Self-service kiosk or receptionist, token display"),
        ("OPD: Consultation", "Full EMR — notes, orders, prescriptions in real-time"),
        ("OPD: Referrals", "Internal specialist referrals, external referral letters"),
        ("OPD: Telemedicine", "Video OPD consultations via integrated platform"),
        ("ER: Rapid Registration", "Minimal info required — full details captured later"),
        ("ER: Triage (ESI/MTS)", "Emergency Severity Index 1-5, Manchester Triage"),
        ("ER: Real-time Bed Board", "Visual display of all ER bays/beds with status"),
        ("ER: STAT Orders", "Priority lab and imaging with STAT flags"),
        ("ER: Mass Casualty Mode", "Simplified documentation for disaster scenarios"),
    ]),
    ("3.14 OPERATING THEATRE & 3.15 BLOOD BANK", [
        ("OT: Surgery Scheduling", "Book OT slots, manage conflicts, priority scheduling"),
        ("OT: WHO Safety Checklist", "Digital Sign-In, Time-Out, Sign-Out workflow"),
        ("OT: Documentation", "Surgical notes, anesthesia record, implant tracking"),
        ("OT: Resource Management", "Equipment, staff, consumable allocation per surgery"),
        ("OT: Surgeon Pref Cards", "Pre-configured instrument and supply lists"),
        ("Blood Bank: Donors", "Registration, screening, donation history"),
        ("Blood Bank: Cross-Match", "Compatibility testing, results, issue units"),
        ("Blood Bank: Transfusion", "Track transfusions, reaction monitoring"),
        ("Blood Bank: Compliance", "Mandatory screening (HIV, HBV, HCV, etc.)"),
    ]),
]

for title, features in module_slides:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
                title, 24, WHITE, True)

    add_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8),
              ["Feature", "Description"], features)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 24: ADMIN MODULES - SECTION TITLE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_layout(slide, 4, "ADMINISTRATIVE MODULES",
                   "Dashboard, Reports, Users, Audit Trail")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 25: ADMIN DASHBOARD & ANALYTICS
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "4.1 ADMIN DASHBOARD & ANALYTICS", 24, WHITE, True)

# KPI boxes row
kpis = [
    ("👥 1,250", "Active Patients", DARK_BLUE),
    ("🛏️ 78%", "Bed Occupancy", DARK_GREEN),
    ("💰 $355K", "Monthly Revenue", ORANGE),
    ("⏱️ 18 min", "Avg Wait Time", MED_BLUE),
    ("📊 4.2/5", "Patient Satisfaction", PURPLE),
]
for i, (val, label, color) in enumerate(kpis):
    left = Inches(0.4 + i * 2.55)
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.4), Inches(2.3), Inches(1.3)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = val
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(10)
    p2.font.color.rgb = RGBColor(220, 220, 255)
    p2.alignment = PP_ALIGN.CENTER

admin_features = [
    ("Executive Dashboard", "Hospital-wide KPIs, revenue, patient count, bed occupancy"),
    ("Department Dashboards", "Department-specific metrics and performance indicators"),
    ("Financial Overview", "Revenue, expenses, outstanding, cash flow trends"),
    ("Quality Indicators", "Infection rates, readmission rates, mortality rates"),
    ("Custom Widgets", "Drag-and-drop dashboard customization"),
    ("Alert Center", "Stock-outs, equipment downtime, pending approvals"),
]
add_table(slide, Inches(0.5), Inches(3.0), Inches(12.3), Inches(4.0),
          ["Feature", "Description"], admin_features)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 26: REPORTS & USER MANAGEMENT
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "4.2 REPORTS, USER ROLES & AUDIT TRAIL", 24, WHITE, True)

# Reports section
add_textbox(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.5),
            "📊 Report Categories", 16, DARK_BLUE, True)

reports = [
    ("Clinical", "Census, diagnosis distribution, outcomes"),
    ("Financial", "Revenue, billing, department P&L"),
    ("Operational", "OT utilization, bed occupancy, ER TAT"),
    ("Custom Builder", "Drag-and-drop with filters & charts"),
    ("Scheduled", "Auto email daily/weekly/monthly"),
]
add_table(slide, Inches(0.5), Inches(1.9), Inches(6.0), Inches(3.0),
          ["Category", "Details"], reports)

# Roles section
add_textbox(slide, Inches(7), Inches(1.3), Inches(6), Inches(0.5),
            "👤 System Roles (RBAC)", 16, DARK_BLUE, True)

roles = [
    ("Administrator", "Full system access"),
    ("Doctor", "Records, prescriptions, orders"),
    ("Nurse", "Vitals, eMAR, care plans"),
    ("Receptionist", "Registration, appointments"),
    ("Pharmacist", "Dispensing, inventory"),
    ("Lab Technician", "Samples, results, QC"),
    ("Accountant", "Billing, claims, reports"),
    ("Custom Roles", "Granular permissions"),
]
add_table(slide, Inches(7), Inches(1.9), Inches(5.8), Inches(4.5),
          ["Role", "Access Scope"], roles)

# Audit trail note
add_shape_bg(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8), LIGHT_BLUE)
add_textbox(slide, Inches(0.7), Inches(5.3), Inches(12), Inches(0.5),
            "🔒 Complete Audit Trail", 14, DARK_BLUE, True)
add_textbox(slide, Inches(0.7), Inches(5.8), Inches(11.5), Inches(1.0),
            "Every action logged: who, what, when, where, before/after values. "
            "Data access logging for patient records. Full version history of all clinical documents. "
            "HIPAA & GDPR compliant with consent management and data portability.", 11, BLACK)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 27: WORKFLOW DIAGRAMS - SECTION TITLE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_layout(slide, 5, "WORKFLOW DIAGRAMS",
                   "How Information Flows Through the System")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 28: PATIENT JOURNEY WORKFLOW
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "COMPLETE PATIENT JOURNEY", 28, WHITE, True)

# Row 1: Arrival flow
row1_steps = ["Patient\nArrives", "Registration\n(New/Existing)", "Triage /\nOPD Queue", "Doctor\nConsultation"]
add_horizontal_flow(slide, row1_steps, Inches(0.5), Inches(1.5), Inches(2.2), Inches(1.0), Inches(0.7), DARK_BLUE)

# Row 2: Action branches
add_textbox(slide, Inches(5.5), Inches(2.7), Inches(2.5), Inches(0.5),
            "▼  Doctor Decides  ▼", 12, MED_BLUE, True, PP_ALIGN.CENTER)

row2_steps = ["Lab\nTests", "Imaging /\nRadiology", "Pharmacy\nMedications", "Admit to\nWard", "Refer\nExternal"]
for i, step in enumerate(row2_steps):
    left = Inches(0.3 + i * 2.6)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3.3), Inches(2.2), Inches(0.9))
    color = [DARK_GREEN, PURPLE, ORANGE, RED, GRAY][i]
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = step
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

# Row 3: Outcomes
add_textbox(slide, Inches(4.0), Inches(4.5), Inches(5), Inches(0.4),
            "▼  Results Return to Doctor  ▼", 12, MED_BLUE, True, PP_ALIGN.CENTER)

row3_steps = ["Follow-up\nScheduled", "Billing &\nPayment", "Discharge\nSummary", "Insurance\nClaim Filed"]
add_horizontal_flow(slide, row3_steps, Inches(1.2), Inches(5.1), Inches(2.4), Inches(1.0), Inches(0.6), ACCENT_TEAL)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 29: BILLING & LAB WORKFLOWS
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "BILLING & LABORATORY WORKFLOWS", 28, WHITE, True)

# Billing workflow (left)
add_textbox(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(0.5),
            "💰 Billing / Revenue Cycle", 16, DARK_BLUE, True)

billing_steps = [
    "Services Rendered", "Charges Captured", "Insurance Check",
    "Invoice Generated", "Payment Collected", "Claim Submitted",
    "Reconciliation", "Receipt Issued"
]
add_flowchart_boxes(slide, billing_steps, Inches(1.5), Inches(1.9),
                    Inches(3.5), Inches(0.4), Inches(0.18), ORANGE)

# Lab workflow (right)
add_textbox(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.5),
            "🔬 Laboratory Workflow", 16, DARK_BLUE, True)

lab_steps = [
    "Doctor Orders Test", "Order in Lab Queue", "Sample Collection",
    "Processing (Analyzer)", "Result Entry / Capture",
    "Validation (QC)", "Report to EMR", "Doctor Notified"
]
add_flowchart_boxes(slide, lab_steps, Inches(8.0), Inches(1.9),
                    Inches(3.5), Inches(0.4), Inches(0.18), DARK_GREEN)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 30: ER & SURGERY WORKFLOWS
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "EMERGENCY & SURGERY WORKFLOWS", 28, WHITE, True)

# ER workflow (left)
add_textbox(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(0.5),
            "🚑 Emergency Department", 16, RED, True)

er_steps = [
    "Ambulance / Walk-in", "Rapid Registration", "Triage (ESI 1-5)",
    "Assign ER Bay/Bed", "STAT Labs & Imaging", "Emergency Treatment",
    "Disposition: Admit/DC/Transfer"
]
add_flowchart_boxes(slide, er_steps, Inches(1.5), Inches(1.9),
                    Inches(3.5), Inches(0.45), Inches(0.2), RED)

# Surgery workflow (right)
add_textbox(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.5),
            "🏥 Operating Theatre", 16, PURPLE, True)

ot_steps = [
    "Surgery Scheduled", "Pre-Op Assessment", "Consent Forms Signed",
    "WHO Checklist: Sign-In", "WHO Checklist: Time-Out", "Surgery Performed",
    "WHO Checklist: Sign-Out", "Recovery / PACU", "Transfer to Ward"
]
add_flowchart_boxes(slide, ot_steps, Inches(8.0), Inches(1.9),
                    Inches(3.5), Inches(0.38), Inches(0.16), PURPLE)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 31: INTEGRATION - SECTION TITLE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_layout(slide, 6, "INTEGRATION & INTEROPERABILITY",
                   "Connect with Everything — Standards-Based")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 32: INTEGRATIONS TABLE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "INTEGRATIONS & STANDARDS", 28, WHITE, True)

integrations = [
    ("HL7 / FHIR", "Healthcare Data Exchange", "Interoperability with other hospital systems"),
    ("DICOM", "Medical Imaging", "Radiology equipment, PACS servers"),
    ("ICD-10 / ICD-11", "Disease Classification", "Standardized diagnosis coding"),
    ("SNOMED CT", "Clinical Terminology", "Structured clinical data recording"),
    ("Lab Instruments", "ASTM / HL7 v2", "Bi-directional analyzer interface"),
    ("Insurance Portals", "Electronic Claims", "Pre-auth, eligibility, adjudication"),
    ("Payment Gateways", "Stripe / PayStack", "Online payments, mobile money"),
    ("Accounting", "QuickBooks / Xero", "Financial data synchronization"),
    ("SMS / Email", "Twilio / SendGrid", "Patient notifications & reminders"),
    ("Video Calling", "Jitsi / Zoom", "Telemedicine consultations"),
    ("Biometric Devices", "Fingerprint / Face", "Staff attendance, patient ID"),
    ("National Health ID", "Government Systems", "Link with national registries"),
]
add_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8),
          ["Standard / System", "Category", "Purpose"], integrations)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 33: SECURITY - SECTION TITLE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_layout(slide, 7, "SECURITY & COMPLIANCE",
                   "Bank-Grade Security for Healthcare Data")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 34: SECURITY DETAILS
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "SECURITY & COMPLIANCE FEATURES", 28, WHITE, True)

security = [
    ("🔒 Data Encryption", "AES-256 at rest, TLS 1.3 in transit — all data encrypted"),
    ("🛡️ Access Control", "Role-based (RBAC), attribute-based (ABAC), IP whitelisting"),
    ("📱 Multi-Factor Auth", "SMS OTP, authenticator app, biometric second factor"),
    ("📋 Audit Logging", "Every action logged: who, what, when, where, before/after"),
    ("💾 Data Backup", "Hourly incremental, daily full, off-site replication"),
    ("🔄 Disaster Recovery", "RPO < 1 hour, RTO < 4 hours, automated failover"),
    ("🏥 HIPAA Compliance", "Full health data protection compliance"),
    ("🌍 GDPR Support", "Right to access, erasure, data portability"),
    ("✅ SOC 2 Type II", "Service organization security certification"),
    ("🔍 Penetration Testing", "Regular third-party security audits"),
    ("🔑 Session Management", "Auto-timeout, concurrent limits, device tracking"),
    ("📍 Data Residency", "Choose data storage location for regulations"),
]
add_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8),
          ["Security Feature", "Details"], security)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 35: TECHNOLOGY STACK - SECTION TITLE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_layout(slide, 8, "TECHNOLOGY STACK",
                   "Modern, Proven, Open-Source Foundation")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 36: TECH STACK VISUAL
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "TECHNOLOGY STACK", 28, WHITE, True)

tech_layers = [
    ("FRONTEND", MED_BLUE, ["React.js", "Next.js 14", "TypeScript", "Tailwind CSS", "Shadcn/UI"]),
    ("MOBILE", ACCENT_TEAL, ["React Native", "Flutter", "Expo", "Offline-First"]),
    ("BACKEND", DARK_BLUE, ["Node.js", "NestJS", "TypeScript", "GraphQL", "REST API", "Socket.IO"]),
    ("DATABASE", DARK_GREEN, ["PostgreSQL", "Redis", "Elasticsearch", "MinIO/S3"]),
    ("DEVOPS", ORANGE, ["Docker", "Kubernetes", "GitHub Actions", "Terraform"]),
    ("MONITORING", PURPLE, ["Prometheus", "Grafana", "ELK Stack", "Sentry"]),
]

for layer_idx, (layer_name, color, techs) in enumerate(tech_layers):
    top = Inches(1.3 + layer_idx * 1.0)
    # Layer label
    label = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), top, Inches(1.8), Inches(0.7)
    )
    label.fill.solid()
    label.fill.fore_color.rgb = color
    label.line.fill.background()
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = layer_name
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Tech items
    for t_idx, tech in enumerate(techs):
        left = Inches(2.3 + t_idx * 1.9)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.7), Inches(0.7)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE
        box.line.color.rgb = color
        box.line.width = Pt(1.5)
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = tech
        p.font.size = Pt(10)
        p.font.color.rgb = color
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 37: IMPLEMENTATION - SECTION TITLE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_layout(slide, 9, "IMPLEMENTATION PLAN",
                   "Phased Delivery Over 12 Months")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 38: IMPLEMENTATION TIMELINE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "12-MONTH IMPLEMENTATION TIMELINE", 28, WHITE, True)

phases = [
    ("Phase 1: Foundation", "Months 1-3", DARK_BLUE,
     "Patient registration, User management, Basic EMR, Appointment scheduling"),
    ("Phase 2: Clinical Core", "Months 4-6", DARK_GREEN,
     "Full EMR + templates, OPD workflow, Lab (LIS), Pharmacy, Billing, Nursing"),
    ("Phase 3: Advanced", "Months 7-9", PURPLE,
     "IPD/Wards, Emergency Dept, Operating Theatre, Radiology/PACS, Blood Bank"),
    ("Phase 4: Analytics", "Months 10-11", ORANGE,
     "Admin dashboard, BI & Reports, Insurance integration, Payment gateways"),
    ("Phase 5: Launch", "Month 12", RED,
     "Performance tuning, Security hardening, Training, Data migration, GO-LIVE"),
]

for i, (name, timeline, color, deliverables) in enumerate(phases):
    top = Inches(1.3 + i * 1.2)
    # Phase bar
    bar_width = Inches(10.5) if i < 3 else (Inches(7.0) if i == 3 else Inches(3.5))
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, Inches(2.5), Inches(1.0)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = timeline
    p2.font.size = Pt(9)
    p2.font.color.rgb = RGBColor(220, 220, 255)
    p2.alignment = PP_ALIGN.CENTER

    # Deliverables
    add_textbox(slide, Inches(3.2), top + Inches(0.1), Inches(9.8), Inches(0.8),
                deliverables, 10, BLACK, False, PP_ALIGN.LEFT)

# Gantt-style bar at bottom
add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
            "M1    M2    M3    M4    M5    M6    M7    M8    M9    M10   M11   M12",
            10, GRAY, True, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 39: PRICING - SECTION TITLE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_layout(slide, 10, "PRICING & LICENSING",
                   "Flexible Plans for Every Hospital Size")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 40: PRICING TABLE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "PRICING & LICENSING OPTIONS", 28, WHITE, True)

pricing = [
    ("Starter", "Up to 50 beds", "Core EMR, OPD, Pharmacy, Lab, Billing", "Custom Quote"),
    ("Professional", "50 – 200 beds", "All Starter + IPD, ER, OT, HR, Analytics", "Custom Quote"),
    ("Enterprise", "200+ beds", "Full suite + PACS, API access, custom dev", "Custom Quote"),
    ("Cloud Hosted", "Any size", "SaaS model — monthly per-user pricing", "Custom Quote"),
    ("On-Premise", "Any size", "Self-hosted with perpetual license", "Custom Quote"),
]
add_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(3.0),
          ["Plan", "Hospital Size", "Includes", "Price"], pricing)

# All plans include
add_shape_bg(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.5), LIGHT_BLUE)
add_textbox(slide, Inches(0.7), Inches(4.7), Inches(12), Inches(0.5),
            "✅ ALL PLANS INCLUDE:", 14, DARK_BLUE, True)

included = [
    "Implementation support & project management",
    "Data migration from existing systems",
    "User training (on-site + virtual sessions)",
    "12-month warranty & bug fixes",
    "Technical documentation & user guides",
    "Dedicated account manager",
]
add_multi_text(slide, Inches(0.7), Inches(5.2), Inches(11.5), Inches(1.8),
               included, 12, BLACK, 1.3, bullet=True)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 41: SUPPORT - SECTION TITLE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_layout(slide, 11, "SUPPORT & MAINTENANCE",
                   "We're With You Every Step of the Way")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 42: SUPPORT DETAILS
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "SUPPORT & MAINTENANCE", 28, WHITE, True)

support = [
    ("24/7 Help Desk", "Round-the-clock support via phone, email, and chat"),
    ("Dedicated Account Mgr", "Single point of contact for your hospital"),
    ("Remote Support", "Screen-sharing for troubleshooting and guidance"),
    ("On-Site Support", "Available for critical issues and installations"),
    ("Regular Updates", "Security patches, feature updates, bug fixes"),
    ("Annual Health Check", "System performance review and optimization"),
    ("Training Sessions", "Quarterly refresher training for staff"),
    ("Knowledge Base", "Online documentation and video tutorials"),
    ("SLA Guarantee", "99.9% uptime, <1hr response for critical issues"),
    ("Data Migration", "Expert assistance migrating from existing systems"),
]
add_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.5),
          ["Service", "Description"], support)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 43: WHY CHOOSE US - SECTION TITLE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_layout(slide, 12, "WHY CHOOSE OUR EMR?",
                   "10 Reasons to Partner with Us")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 44: WHY CHOOSE US - CONTENT
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "10 REASONS TO CHOOSE OUR EMR", 28, WHITE, True)

reasons = [
    ("Purpose-Built", "Designed specifically for private hospitals — not generic"),
    ("All-in-One", "No need for multiple vendor systems — everything integrated"),
    ("Modern Tech", "Fast, responsive, mobile-ready — built with latest tools"),
    ("Scalable", "Grows with you from 10-bed clinic to 1000+ bed hospital"),
    ("Interoperable", "HL7/FHIR/DICOM compliant — integrates with anything"),
    ("Secure", "Bank-grade AES-256 encryption + full audit trail"),
    ("Affordable", "Flexible licensing, no hidden costs, strong ROI"),
    ("Local Support", "Dedicated support team in your region"),
    ("Fast Launch", "Go live in as little as 3 months (Phase 1 MVP)"),
    ("Proven ROI", "Hospitals report 30-50% efficiency improvement"),
]

for i, (title, desc) in enumerate(reasons):
    col = i % 2
    row = i // 2
    left = Inches(0.5 + col * 6.4)
    top = Inches(1.3 + row * 1.2)
    # Number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, Inches(0.55), Inches(0.55)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = DARK_BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_textbox(slide, left + Inches(0.7), top, Inches(5.3), Inches(0.35),
                title, 14, DARK_BLUE, True)
    add_textbox(slide, left + Inches(0.7), top + Inches(0.35), Inches(5.3), Inches(0.5),
                desc, 11, GRAY)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 45: IN-PATIENT WORKFLOW
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "IN-PATIENT & PHARMACY WORKFLOWS", 28, WHITE, True)

# IPD workflow (left)
add_textbox(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(0.5),
            "🛏️ In-Patient Workflow", 16, DARK_BLUE, True)

ipd_steps = [
    "Admission Request", "Check Bed Availability", "Assign Bed (Ward/ICU)",
    "Admission Documentation", "Daily Rounds & Orders", "Nursing Care",
    "Discharge Planning", "All-Dept Clearance", "Discharge Summary"
]
add_flowchart_boxes(slide, ipd_steps, Inches(1.5), Inches(1.9),
                    Inches(3.5), Inches(0.38), Inches(0.15), DARK_BLUE)

# Pharmacy workflow (right)
add_textbox(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.5),
            "💊 Pharmacy Workflow", 16, DARK_GREEN, True)

pharm_steps = [
    "E-Prescription Received", "Verify Patient & Allergy",
    "Drug Interaction Check", "Dispense Medication",
    "Barcode Scan & Verify", "Label Printing",
    "Patient Counselling", "Billing & Payment", "Update Inventory"
]
add_flowchart_boxes(slide, pharm_steps, Inches(8.0), Inches(1.9),
                    Inches(3.5), Inches(0.38), Inches(0.15), DARK_GREEN)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 46: SAMPLE DASHBOARD PREVIEW
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "SAMPLE DASHBOARD — KPI OVERVIEW", 28, WHITE, True)

# KPI Cards row 1
kpi_cards = [
    ("Total Patients\nToday", "1,247", "▲ 12%", DARK_BLUE),
    ("Bed Occupancy\nRate", "78%", "▲ 3%", DARK_GREEN),
    ("Revenue\nThis Month", "$355,400", "▲ 8%", ORANGE),
    ("Average Wait\nTime (OPD)", "18 min", "▼ 22%", MED_BLUE),
    ("Lab TAT\n(Routine)", "2.5 hrs", "▼ 15%", PURPLE),
]
for i, (label, value, trend, color) in enumerate(kpi_cards):
    left = Inches(0.3 + i * 2.55)
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.3), Inches(2.4), Inches(1.8)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(2)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.size = Pt(22)
    p2.font.color.rgb = color
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.text = trend
    p3.font.size = Pt(10)
    p3.font.color.rgb = DARK_GREEN if "▲" in trend else RED
    p3.alignment = PP_ALIGN.CENTER

# Department stats table
dept_stats = [
    ("OPD", "450", "18 min", "$45,200", "4.3/5"),
    ("Emergency", "80", "8 min", "$38,100", "4.1/5"),
    ("Surgery", "12", "On schedule", "$85,000", "4.5/5"),
    ("Laboratory", "680 tests", "2.5 hrs TAT", "$32,400", "4.2/5"),
    ("Pharmacy", "890 Rx", "12 min avg", "$55,300", "4.0/5"),
    ("Radiology", "120 scans", "45 min TAT", "$28,600", "4.4/5"),
    ("IPD (Wards)", "96 beds used", "Avg 3.2 days", "$72,000", "4.3/5"),
]
add_table(slide, Inches(0.3), Inches(3.4), Inches(12.7), Inches(3.8),
          ["Department", "Volume", "Efficiency", "Revenue", "Satisfaction"],
          dept_stats)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 47: COMPETITIVE ADVANTAGES
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "COMPETITIVE COMPARISON", 28, WHITE, True)

comp = [
    ("Feature", "Our EMR", "Competitor A", "Competitor B"),
    ("All Modules Integrated", "✅ Yes", "❌ Separate vendors", "⚠️ Partial"),
    ("Cloud + On-Premise", "✅ Both options", "☁️ Cloud only", "🏢 On-prem only"),
    ("Mobile App", "✅ Full-featured", "❌ None", "⚠️ Basic"),
    ("Telemedicine Built-in", "✅ Integrated", "❌ 3rd party", "❌ No"),
    ("Custom Report Builder", "✅ Drag & drop", "⚠️ Limited", "⚠️ Limited"),
    ("PACS/DICOM", "✅ Integrated", "❌ Separate", "❌ No"),
    ("AI / Clinical Decision", "✅ Included", "❌ No", "❌ No"),
    ("HL7/FHIR Compliant", "✅ Full", "⚠️ Partial", "❌ No"),
    ("Implementation Time", "3-12 months", "12-18 months", "6-12 months"),
    ("Total Cost of Ownership", "💰 Lower", "💸 Higher", "💰 Medium"),
]

# Custom table for comparison
table_shape = slide.shapes.add_table(len(comp), 4, Inches(0.5), Inches(1.3),
                                      Inches(12.3), Inches(5.8))
table = table_shape.table
for ci in range(4):
    table.columns[ci].width = Inches(3.08)

for r_idx, row in enumerate(comp):
    for c_idx, val in enumerate(row):
        cell = table.cell(r_idx, c_idx)
        cell.text = val
        is_header = (r_idx == 0)
        if is_header:
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE
        elif c_idx == 1:  # Our EMR column
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GREEN
        elif r_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.color.rgb = WHITE if is_header else BLACK
            p.font.bold = is_header or (c_idx == 0)
            p.font.name = "Calibri"
            if c_idx > 0:
                p.alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 48: NEXT STEPS
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            "RECOMMENDED NEXT STEPS", 28, WHITE, True)

next_steps = [
    ("1", "Discovery Meeting", "Deep-dive into your hospital's specific requirements,\ndepartments, workflows, and pain points", MED_BLUE),
    ("2", "Live Demo", "See the system in action with your real-world\nscenarios and clinical workflows", DARK_GREEN),
    ("3", "Proposal & Scope", "Detailed proposal with scope, timeline,\ncost, and implementation plan", ORANGE),
    ("4", "Pilot Program", "Start with Phase 1 (3 months) in one department\nto prove value before full rollout", PURPLE),
    ("5", "Full Implementation", "Phased rollout across all departments\nover 12 months with training & support", RED),
]

for i, (num, title, desc, color) in enumerate(next_steps):
    left = Inches(0.3 + i * 2.55)
    top = Inches(1.5)

    # Number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.6), top, Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Title
    add_textbox(slide, left, top + Inches(1.0), Inches(2.4), Inches(0.5),
                title, 14, color, True, PP_ALIGN.CENTER)

    # Description
    add_textbox(slide, left, top + Inches(1.5), Inches(2.4), Inches(1.5),
                desc, 10, BLACK, False, PP_ALIGN.CENTER)

# Bottom CTA
add_shape_bg(slide, Inches(2), Inches(5.5), Inches(9.3), Inches(1.5), DARK_BLUE)
add_textbox(slide, Inches(2.2), Inches(5.6), Inches(9), Inches(0.6),
            "Let's schedule a discovery meeting to understand your needs",
            16, WHITE, True, PP_ALIGN.CENTER)
add_textbox(slide, Inches(2.2), Inches(6.2), Inches(9), Inches(0.5),
            "We'll prepare a customized demo based on your hospital's workflow",
            12, RGBColor(180, 210, 240), False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 49: CLOSING / THANK YOU
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)

add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), MED_BLUE)

add_textbox(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1.2),
            "Ready to Transform", 44, WHITE, True, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(2.6), Inches(11.3), Inches(1.2),
            "Your Hospital?", 44, RGBColor(100, 180, 255), True, PP_ALIGN.CENTER)

add_textbox(slide, Inches(2), Inches(4.0), Inches(9.3), Inches(0.7),
            "Contact us for a personalized demo and consultation",
            18, RGBColor(180, 210, 240), False, PP_ALIGN.CENTER)

# Contact info boxes
contact_items = [
    ("📧 Email", "info@hospitalemr.com"),
    ("📞 Phone", "+XXX-XXX-XXXX"),
    ("🌐 Website", "www.hospitalemr.com"),
]
for i, (label, value) in enumerate(contact_items):
    left = Inches(2.5 + i * 3.0)
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(5.0), Inches(2.5), Inches(1.0)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0, 40, 80)
    box.line.color.rgb = MED_BLUE
    box.line.width = Pt(1)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(150, 180, 210)
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.size = Pt(14)
    p2.font.color.rgb = WHITE
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(6)

add_textbox(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.5),
            "CONFIDENTIAL  |  © 2026 Hospital EMR Solutions  |  All Rights Reserved",
            10, RGBColor(120, 140, 160), False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════
output_path = os.path.join(OUTPUT_DIR, "Hospital_EMR_Pitch.pptx")
prs.save(output_path)
print(f"✅ PowerPoint saved to: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
